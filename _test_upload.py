"""
Validation script for the EduLingo upload module.
Run with:  python _test_upload.py
"""
import io
import sys

import app as a

client = a.app.test_client()

errors = []

def check(label, condition, detail=""):
    if condition:
        print(f"  OK  {label}")
    else:
        msg = f"FAIL  {label}" + (f": {detail}" if detail else "")
        print(f"  {msg}")
        errors.append(msg)

# ── 1. Page routes ────────────────────────────────────────────────────────
print("-- Page routes --")
for route in ["/", "/translate", "/summarize", "/explain",
              "/quiz", "/ask", "/chat", "/history"]:
    code = client.get(route).status_code
    check(f"GET {route:<12} -> {code}", code == 200)

# ── 2. API validation guards ──────────────────────────────────────────────
print("\n-- API guards --")

resp = client.post("/api/upload", data={}, content_type="multipart/form-data")
check("No file field -> 400", resp.status_code == 400)

resp = client.post(
    "/api/upload",
    data={"file": (io.BytesIO(b"data"), "test.xyz")},
    content_type="multipart/form-data",
)
j = resp.get_json()
check("Bad extension -> 415 + success=false", resp.status_code == 415 and not j["success"])

# ── 3. TXT upload ─────────────────────────────────────────────────────────
print("\n-- File type uploads --")

txt_content = b"Hello world. This is a test document with some words in it."
resp = client.post(
    "/api/upload",
    data={"file": (io.BytesIO(txt_content), "sample.txt")},
    content_type="multipart/form-data",
)
j = resp.get_json()
check(
    "TXT upload  words={wc}  session={sid}".format(
        wc=j.get("word_count"), sid=str(j.get("session_id", ""))[:8]
    ),
    j.get("success") and j.get("file_type") == "TXT",
    j,
)

# ── 4. DOCX upload ────────────────────────────────────────────────────────
from docx import Document

buf = io.BytesIO()
doc = Document()
doc.add_paragraph("EduLingo DOCX test paragraph. Contains multiple words here.")
doc.save(buf)
buf.seek(0)

resp = client.post(
    "/api/upload",
    data={"file": (buf, "sample.docx")},
    content_type="multipart/form-data",
)
j = resp.get_json()
check(
    "DOCX upload  words={wc}  session={sid}".format(
        wc=j.get("word_count"), sid=str(j.get("session_id", ""))[:8]
    ),
    j.get("success") and j.get("file_type") == "DOCX",
    j,
)

# ── 5. PPTX upload ────────────────────────────────────────────────────────
from pptx import Presentation

prs   = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text      = "EduLingo PPTX Test"
slide.placeholders[1].text   = "This slide contains test content for the upload module."
buf2  = io.BytesIO()
prs.save(buf2)
buf2.seek(0)

resp = client.post(
    "/api/upload",
    data={"file": (buf2, "sample.pptx")},
    content_type="multipart/form-data",
)
j = resp.get_json()
check(
    "PPTX upload  slides={pg}  words={wc}  session={sid}".format(
        pg=j.get("pages"), wc=j.get("word_count"), sid=str(j.get("session_id", ""))[:8]
    ),
    j.get("success") and j.get("file_type") == "PPTX",
    j,
)

# ── 6. PDF upload ─────────────────────────────────────────────────────────
from PyPDF2 import PdfWriter

writer = PdfWriter()
writer.add_blank_page(width=612, height=792)
buf3 = io.BytesIO()
writer.write(buf3)
buf3.seek(0)

resp = client.post(
    "/api/upload",
    data={"file": (buf3, "sample.pdf")},
    content_type="multipart/form-data",
)
j = resp.get_json()
check(
    "PDF  upload  pages={pg}  words={wc}  session={sid}".format(
        pg=j.get("pages"), wc=j.get("word_count"), sid=str(j.get("session_id", ""))[:8]
    ),
    j.get("success") and j.get("file_type") == "PDF",
    j,
)

# ── 7. Document store populated ───────────────────────────────────────────
print("\n-- Document store --")
count = len(a.document_store)
check("document_store has >= 4 entries (got {n})".format(n=count), count >= 4)

# ── Summary ───────────────────────────────────────────────────────────────
print()
if errors:
    print("FAILED — {n} error(s):".format(n=len(errors)))
    for e in errors:
        print(" ", e)
    sys.exit(1)
else:
    print("All checks passed.")
