import os
import random
import re
import uuid
import datetime

from dotenv import load_dotenv

_ENV_PATH = os.path.join(os.path.abspath(os.path.dirname(__file__)), "uploads", ".env")
load_dotenv(_ENV_PATH)
WATSONX_API_KEY = os.environ.get("WATSONX_API_KEY", "")
WATSONX_PROJECT_ID = os.environ.get("WATSONX_PROJECT_ID", "")
WATSONX_URL = os.environ.get("WATSONX_URL", "https://us-south.ml.cloud.ibm.com")
WATSONX_MODEL_ID = os.environ.get("WATSONX_MODEL_ID", "meta-llama/llama-3-3-70b-instruct")

from flask import Flask, jsonify, render_template, request

# ── Optional parsing libraries ─────────────────────────────────────────────
try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    from pptx import Presentation as PptxPresentation
except ImportError:
    PptxPresentation = None

# ── Optional IBM watsonx.ai SDK ────────────────────────────────────────────
try:
    from ibm_watsonx_ai import Credentials
    from ibm_watsonx_ai.foundation_models import ModelInference
    from ibm_watsonx_ai.foundation_models.schema import TextGenParameters
    WATSONX_AVAILABLE = True
except ImportError:
    WATSONX_AVAILABLE = False

# ── App setup ──────────────────────────────────────────────────────────────
app = Flask(__name__)

BASE_DIR = os.path.abspath(os.path.dirname(__file__))
app.config["UPLOAD_FOLDER"]     = os.path.join(BASE_DIR, "uploads")
app.config["TRANSLATED_FOLDER"] = os.path.join(BASE_DIR, "translated")
app.config["MAX_CONTENT_LENGTH"] = 50 * 1024 * 1024   # 50 MB

os.makedirs(app.config["UPLOAD_FOLDER"],     exist_ok=True)
os.makedirs(app.config["TRANSLATED_FOLDER"], exist_ok=True)

# ── Paths for durable stores ───────────────────────────────────────────────
_DOCUMENT_STORE_PATH = os.path.join(BASE_DIR, "uploads", ".document_store.json")

# ── In-memory stores ───────────────────────────────────────────────────────
# { session_id: { text, filename, file_type, pages, word_count } }
document_store: dict = {}

# [ { id, type, session_id, filename, target_language, word_count,
#     mode, timestamp, preview } ]
history_store: list = []

# { quiz_id: { session_id, filename, questions, mode, created_at } }
quiz_store: dict = {}

ALLOWED_EXTENSIONS = {"pdf", "docx", "pptx", "txt"}


# ── Document-store persistence ─────────────────────────────────────────────

def _load_document_store() -> None:
    """
    Populate document_store from the JSON sidecar on disk.
    Called once at startup so sessions survive a Flask reload/restart.
    Only entries whose source file still exists on disk are loaded.
    """
    import json as _json
    if not os.path.isfile(_DOCUMENT_STORE_PATH):
        return
    try:
        with open(_DOCUMENT_STORE_PATH, "r", encoding="utf-8") as fh:
            saved = _json.load(fh)
        for sid, entry in saved.items():
            ext       = entry.get("file_type", "").lower()
            src_file  = os.path.join(app.config["UPLOAD_FOLDER"], f"{sid}.{ext}")
            if os.path.isfile(src_file):
                document_store[sid] = entry
    except Exception:
        pass   # corrupt file — start clean


def _flush_document_store() -> None:
    """Write document_store to the JSON sidecar (called after every upload)."""
    import json as _json
    try:
        with open(_DOCUMENT_STORE_PATH, "w", encoding="utf-8") as fh:
            _json.dump(document_store, fh, ensure_ascii=False)
    except Exception:
        pass


# Load persisted sessions immediately so existing session_ids are valid
# the moment the server starts (or reloads).
_load_document_store()

SUPPORTED_LANGUAGES = {
    "English":   "English",
    "Hindi":     "Hindi",
    "Telugu":    "Telugu",
    "Tamil":     "Tamil",
    "Kannada":   "Kannada",
    "Malayalam": "Malayalam",
    "Marathi":   "Marathi",
    "Bengali":   "Bengali",
}

# ── watsonx.ai configuration (read from environment) ──────────────────────
WATSONX_API_KEY    = os.environ.get("WATSONX_API_KEY", "")
WATSONX_PROJECT_ID = os.environ.get("WATSONX_PROJECT_ID", "")
WATSONX_URL        = os.environ.get("WATSONX_URL", "https://us-south.ml.cloud.ibm.com")
WATSONX_MODEL_ID   = os.environ.get("WATSONX_MODEL_ID", "meta-llama/llama-3-3-70b-instruct")

def _mask_secret(value: str) -> str:
    if not value:
        return "<empty>"
    if len(value) <= 4:
        return "*" * len(value)
    return "*" * (len(value) - 4) + value[-4:]


print("API_KEY:", "Loaded" if WATSONX_API_KEY else "Not Loaded", _mask_secret(WATSONX_API_KEY))
print("PROJECT_ID:", "Loaded" if WATSONX_PROJECT_ID else "Not Loaded")
print("URL:", "Loaded" if WATSONX_URL else "Not Loaded")

_QUIZ_STOPWORDS = {
    "about", "after", "again", "also", "because", "been", "before", "being",
    "between", "both", "could", "during", "each", "from", "have", "having",
    "here", "into", "many", "more", "most", "only", "other", "some",
    "such", "than", "that", "them", "then", "there", "these", "they",
    "this", "those", "through", "under", "using", "very", "when", "where",
    "which", "while", "with", "without", "would", "your", "their", "will",
    "shall", "onto", "over", "under", "what", "who", "whom", "why", "how",
    "the", "and", "for", "are", "was", "were", "has", "had", "have",
    "not", "but", "can", "may", "our", "out", "use", "used", "useful",
}


# ── Helpers ────────────────────────────────────────────────────────────────

def allowed_extension(filename: str) -> bool:
    """Return True when the filename has a permitted extension."""
    return (
        "." in filename
        and filename.rsplit(".", 1)[1].lower() in ALLOWED_EXTENSIONS
    )


def extract_text(filepath: str, ext: str) -> tuple[str, int]:
    """
    Extract all readable text from *filepath*.
    Returns (text, page_count).
    """
    ext = ext.lower()

    if ext == "pdf":
        if PyPDF2 is None:
            raise RuntimeError("PyPDF2 is not installed.")
        with open(filepath, "rb") as fh:
            reader = PyPDF2.PdfReader(fh)
            pages  = len(reader.pages)
            text   = "\n".join(
                (page.extract_text() or "") for page in reader.pages
            )
        return text, pages

    if ext == "docx":
        if DocxDocument is None:
            raise RuntimeError("python-docx is not installed.")
        doc  = DocxDocument(filepath)
        text = "\n".join(p.text for p in doc.paragraphs)
        return text, 0

    if ext == "pptx":
        if PptxPresentation is None:
            raise RuntimeError("python-pptx is not installed.")
        prs   = PptxPresentation(filepath)
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    parts.append(shape.text_frame.text)
        text  = "\n".join(parts)
        pages = len(prs.slides)
        return text, pages

    if ext == "txt":
        with open(filepath, "r", encoding="utf-8", errors="replace") as fh:
            text = fh.read()
        return text, 0

    raise ValueError(f"Unsupported extension: {ext}")


def _hard_split_by_words(text: str, max_words: int) -> list[str]:
    """Split *text* into fixed-size word windows when no softer boundary exists."""
    words = text.split()
    if not words:
        return []
    return [
        " ".join(words[i : i + max_words])
        for i in range(0, len(words), max_words)
    ]


def _extract_model_text(response) -> str:
    """Normalise watsonx.ai generate_text output to a plain string."""
    if response is None:
        return ""
    if isinstance(response, str):
        return response.strip()
    if isinstance(response, dict):
        for key in ("generated_text", "text", "result"):
            if response.get(key):
                return str(response[key]).strip()
        results = response.get("results")
        if isinstance(results, list) and results:
            # Concatenate all results instead of only the first one
            parts = []
            for item in results:
                if isinstance(item, dict) and item.get("generated_text"):
                    parts.append(str(item["generated_text"]).strip())
                else:
                    parts.append(str(item).strip())
            return "\n\n".join(part for part in parts if part)
    if isinstance(response, list):
        parts = [_extract_model_text(item) for item in response]
        return "\n\n".join(part for part in parts if part)
    return str(response).strip()


def _split_into_chunks(text: str, max_words: int = 800) -> list[str]:
    """
    Split *text* into chunks of at most *max_words* words, breaking on
    paragraph boundaries (double-newline, then single-newline) and finally
    on sentence or word boundaries so no content is skipped.

    Logging (printed to the Flask console):
      - Total extracted words
      - Number of chunks created
      - Words in each chunk
    """
    normalised = text.replace("\r\n", "\n").replace("\r", "\n")

    # Prefer paragraph breaks; fall back to single newlines for PDF-style text.
    paragraphs = re.split(r"\n{2,}", normalised)
    paragraphs = [p.strip() for p in paragraphs if p.strip()]
    if len(paragraphs) == 1 and len(paragraphs[0].split()) > max_words:
        line_parts = [ln.strip() for ln in paragraphs[0].split("\n") if ln.strip()]
        if len(line_parts) > 1:
            paragraphs = line_parts

    total_words = len(text.split())
    chunks: list[str] = []
    current_parts: list[str] = []
    current_words = 0

    def _flush_current() -> None:
        nonlocal current_parts, current_words
        if current_parts:
            chunks.append("\n\n".join(current_parts))
            current_parts = []
            current_words = 0

    def _add_piece(piece: str) -> None:
        nonlocal current_words
        piece_words = len(piece.split())
        if piece_words > max_words:
            _flush_current()
            for part in _hard_split_by_words(piece, max_words):
                chunks.append(part)
            return
        if current_words + piece_words > max_words and current_parts:
            _flush_current()
        current_parts.append(piece)
        current_words += piece_words

    for para in paragraphs:
        para_words = len(para.split())

        if para_words > max_words:
            sentences = re.split(r"(?<=[.!?])\s+", para)
            if len(sentences) <= 1:
                for part in _hard_split_by_words(para, max_words):
                    chunks.append(part)
            else:
                for sentence in sentences:
                    if sentence.strip():
                        _add_piece(sentence.strip())
            continue

        _add_piece(para)

    _flush_current()

    if not chunks and text.strip():
        chunks = _hard_split_by_words(text.strip(), max_words)

    # ── Logging ────────────────────────────────────────────────────────────
    print(f"[CHUNK] Total words in document : {total_words}", flush=True)
    print(f"[CHUNK] Number of chunks created: {len(chunks)}", flush=True)
    for i, chunk in enumerate(chunks, 1):
        print(f"[CHUNK]   chunk {i}: {len(chunk.split())} words", flush=True)

    return chunks


def _merge_chunk_texts_iteratively(
    model,
    chunk_texts: list[str],
    *,
    label: str,
    merge_instruction: str,
) -> str:
    """
    Fold a list of per-chunk model outputs into one combined string.
    Pairwise merging keeps every chunk's content in scope for the model.
    """
    if not chunk_texts:
        return ""
    combined = chunk_texts[0]
    for i, nxt in enumerate(chunk_texts[1:], 2):
        print(f"[{label}] merging response {i}/{len(chunk_texts)}", flush=True)
        merge_prompt = (
            f"{merge_instruction}\n\n"
            f"Part A:\n{combined}\n\n"
            f"Part B:\n{nxt}\n\n"
            f"Combined output:"
        )
        params = TextGenParameters(max_new_tokens=1000)
        combined = _extract_model_text(model.generate_text(prompt=merge_prompt, params=params))
    return combined

def _watsonx_translate(text: str, target_language: str) -> str:
    """
    Translate *text* into *target_language* using IBM watsonx.ai.
    Splits on paragraph boundaries so every part of the document is translated.
    Raises RuntimeError on any failure so the caller can fall back.
    """
    if not WATSONX_AVAILABLE:
        raise RuntimeError("ibm-watsonx-ai SDK not installed.")
    if not WATSONX_API_KEY or not WATSONX_PROJECT_ID:
        raise RuntimeError("WATSONX_API_KEY or WATSONX_PROJECT_ID not set.")

    credentials = Credentials(
        url=WATSONX_URL,
        api_key=WATSONX_API_KEY,
    )
    model = ModelInference(
        model_id=WATSONX_MODEL_ID,
        credentials=credentials,
        project_id=WATSONX_PROJECT_ID,
    )

    chunks = _split_into_chunks(text, max_words=600)
    if not chunks:
        return ""

    translated_chunks: list[str] = []

    for i, chunk in enumerate(chunks, 1):
        print(f"[TRANSLATE] processing chunk {i}/{len(chunks)} "
              f"({len(chunk.split())} words)", flush=True)
        prompt = (
            f"You are a professional translator. Your only job is to translate text.\n"
            f"Translate the ENTIRE text below into {target_language}.\n"
            f"Rules:\n"
            f"- Translate EVERY paragraph, sentence, and section from start to finish.\n"
            f"- Do NOT stop after the first paragraph — continue until the last word is translated.\n"
            f"- Preserve the original paragraph structure exactly: keep the same number of paragraphs, "
            f"the same line breaks, and the same blank lines between paragraphs.\n"
            f"- Write the entire output in {target_language} script and language.\n"
            f"- Do NOT write any English words except proper nouns (names, places, brands).\n"
            f"- Do NOT summarize, shorten, omit, paraphrase, explain, or add any commentary.\n"
            f"- Do NOT skip headings, lists, tables, captions, or footnotes.\n"
            f"- Do NOT include the original English text.\n"
            f"- Output ONLY the complete {target_language} translation, nothing else.\n\n"
            f"Text to translate:\n{chunk}\n\n"
            f"{target_language} translation:"
        )
        params = TextGenParameters(max_new_tokens=1000)
        raw_response = model.generate_text(prompt=prompt, params=params)
        response = _extract_model_text(raw_response)
        translated_chunks.append(response)

    print(f"[TRANSLATE] collected {len(translated_chunks)}/{len(chunks)} chunk responses",
          flush=True)
    return "\n\n".join(translated_chunks)


def _demo_translate(text: str, target_language: str) -> str:
    """
    Demo mode translation: returns a clearly labelled placeholder that
    echoes the original text so the full UI pipeline can be tested without
    live credentials.
    """
    preview = text[:200].strip()
    return (
        f"[DEMO MODE — {target_language} translation]\n\n"
        f"IBM watsonx.ai credentials are not configured. "
        f"To enable real translations set the environment variables:\n"
        f"  WATSONX_API_KEY, WATSONX_PROJECT_ID, WATSONX_URL\n\n"
        f"Original text preview:\n{preview}"
        + (" …" if len(text) > 200 else "")
    )


def _watsonx_summarize(text: str) -> tuple[str, list[str]]:
    """
    Summarise *text* with IBM watsonx.ai.
    Splits the document into chunks, summarises each chunk, then merges
    all chunk summaries into one coherent final summary with key points.
    Returns (summary, key_points_list).
    Raises RuntimeError on any failure so the caller can fall back.
    """
    if not WATSONX_AVAILABLE:
        raise RuntimeError("ibm-watsonx-ai SDK not installed.")
    if not WATSONX_API_KEY or not WATSONX_PROJECT_ID:
        raise RuntimeError("WATSONX_API_KEY or WATSONX_PROJECT_ID not set.")

    credentials = Credentials(url=WATSONX_URL, api_key=WATSONX_API_KEY)
    model = ModelInference(
        model_id=WATSONX_MODEL_ID,
        credentials=credentials,
        project_id=WATSONX_PROJECT_ID,
    )

    chunks = _split_into_chunks(text, max_words=800)
    if not chunks:
        return "", []

    # ── Step 1: summarise each chunk independently ─────────────────────────
    chunk_summaries: list[str] = []
    for i, chunk in enumerate(chunks, 1):
        print(f"[SUMMARIZE] processing chunk {i}/{len(chunks)} "
              f"({len(chunk.split())} words)", flush=True)
        prompt = (
            "You are an expert academic summarizer.\n"
            "Read the document excerpt below and write a thorough section summary "
            "that covers ALL information in this excerpt.\n"
            "Rules:\n"
            "- Include every topic, subtopic, argument, definition, example, and finding.\n"
            "- Do NOT omit any section, heading, or detail from this excerpt.\n"
            "- Write 5 to 10 complete sentences in clear, plain English.\n"
            "- Do NOT copy sentences verbatim.\n"
            "- Do NOT add information that is not in the excerpt.\n\n"
            f"Excerpt:\n{chunk}\n\nSummary:"
        )
        params = TextGenParameters(max_new_tokens=1000)
        raw_response = model.generate_text(prompt=prompt, params=params)
        response = _extract_model_text(raw_response)
        chunk_summaries.append(response)

    print(f"[SUMMARIZE] collected {len(chunk_summaries)}/{len(chunks)} chunk summaries",
          flush=True)

    # ── Step 2: merge ALL chunk summaries, then produce the final structure ─
    merged_input = "\n\n".join(chunk_summaries)
    combined_summaries = _merge_chunk_texts_iteratively(
        model,
        chunk_summaries,
        label="SUMMARIZE",
        merge_instruction=(
            "You are an expert academic summarizer. Combine Part A and Part B below "
            "into one unified section summary that preserves ALL information from BOTH "
            "parts. Do not omit any topic, detail, or finding."
        ),
    )
    final_prompt = (
        "You are an expert academic summarizer.\n"
        "Below are section-by-section summaries of a full document.\n"
        "Combine them into one comprehensive, well-structured final summary "
        "of 300 to 600 words using ALL of the following labeled sections:\n\n"
        "Title: <document title or topic>\n"
        "Objective: <what the document aims to do or prove>\n"
        "Introduction: <background, context, and scope of the document>\n"
        "Main Concepts: <key ideas, theories, and definitions covered>\n"
        "Methodology: <approaches, methods, or processes described>\n"
        "Results: <findings, outcomes, evidence, or key information>\n"
        "Conclusion: <conclusions, implications, or recommendations>\n\n"
        "Rules:\n"
        "- Cover ALL information from every section summary — do not omit any topic or detail.\n"
        "- Every section listed above MUST appear in your output.\n"
        "- Each section must contain multiple complete sentences (at least 3 sentences per section).\n"
        "- The total summary must be 300 to 600 words.\n"
        "- Write in clear, plain English.\n"
        "- Do NOT copy sentences verbatim.\n"
        "- Do NOT add any preamble or closing remarks outside the structure.\n\n"
        f"Section summaries:\n{combined_summaries}\n\nFinal summary:"
    )
    params = TextGenParameters(max_new_tokens=1000)
    summary = _extract_model_text(model.generate_text(prompt=final_prompt, params=params))

    # ── Step 3: extract key points from ALL chunk summaries ────────────────
    kp_prompt = (
        "You are an expert academic analyst.\n"
        "Read the section summaries below and extract 5 to 10 key points "
        "that capture the most important ideas, findings, and conclusions "
        "from the ENTIRE document.\n\n"
        "Rules:\n"
        "- Each key point must be a complete, informative sentence.\n"
        "- Cover different parts of the document.\n"
        "- Format each key point as a plain line starting with '- '.\n"
        "- Output ONLY the bullet lines, nothing else.\n\n"
        f"Section summaries:\n{merged_input}\n\nKey points:"
    )
    params = TextGenParameters(max_new_tokens=500)
    kp_raw = _extract_model_text(model.generate_text(prompt=kp_prompt, params=params))
    key_points = [
        line.lstrip("-• ").strip()
        for line in kp_raw.splitlines()
        if line.strip() and not line.strip().startswith("#")
    ][:10]

    return summary, key_points


def _demo_summarize(text: str) -> tuple[str, list[str]]:
    """
    Demo mode summarisation — deterministic, no AI call.
    Produces a brief extractive summary and up to 5 key-point sentences.
    """
    sentences = [s.strip() for s in text.replace("\n", " ").split(".") if len(s.strip()) > 20]
    summary_sentences = sentences[:3] if len(sentences) >= 3 else sentences
    summary = (
        "[DEMO MODE — IBM watsonx.ai credentials are not configured. "
        "Set WATSONX_API_KEY, WATSONX_PROJECT_ID and WATSONX_URL for real AI summaries.]\n\n"
        + ". ".join(summary_sentences) + ("." if summary_sentences else "")
    )
    key_points = [s[:120] for s in sentences[3:8]] if len(sentences) > 3 else [s[:120] for s in sentences]
    if not key_points:
        key_points = ["Document uploaded successfully — no extractable sentences found."]
    return summary, key_points[:5]


def _watsonx_explain(text: str) -> list[dict]:
    """
    Generate a chapter-wise explanation of *text* using IBM watsonx.ai.
    Splits the document into chunks and explains each chunk, then merges
    all sections into one complete explanation.
    Returns a list of section dicts: [{ heading, paragraph, bullets }, …]
    Raises RuntimeError on any failure so the caller can fall back.
    """
    if not WATSONX_AVAILABLE:
        raise RuntimeError("ibm-watsonx-ai SDK not installed.")
    if not WATSONX_API_KEY or not WATSONX_PROJECT_ID:
        raise RuntimeError("WATSONX_API_KEY or WATSONX_PROJECT_ID not set.")

    credentials = Credentials(url=WATSONX_URL, api_key=WATSONX_API_KEY)
    model = ModelInference(
        model_id=WATSONX_MODEL_ID,
        credentials=credentials,
        project_id=WATSONX_PROJECT_ID,
    )

    chunks = _split_into_chunks(text, max_words=800)
    if not chunks:
        return [{"heading": "Explanation", "paragraph": "No content could be extracted.", "bullets": []}]

    all_sections: list[dict] = []

    for i, chunk in enumerate(chunks, 1):
        print(f"[EXPLAIN] processing chunk {i}/{len(chunks)} "
              f"({len(chunk.split())} words)", flush=True)
        prompt = (
            "You are an expert educational tutor.\n"
            "Read the document excerpt below and produce a complete structured explanation "
            "that covers EVERY topic, subtopic, and idea in this excerpt.\n\n"
            "For each section output EXACTLY this format:\n"
            "SECTION: <descriptive section heading>\n"
            "PARAGRAPH: <a detailed paragraph of 5 to 8 sentences explaining this part "
            "in simple, clear English. Explain technical terms. Do not copy sentences verbatim.>\n"
            "BULLETS:\n"
            "- <key point from this section>\n"
            "- <key point from this section>\n"
            "- <key point from this section>\n"
            "- <key point from this section>\n\n"
            "Rules:\n"
            "- Explain EVERY topic from the excerpt — do not skip or omit any content.\n"
            "- Produce multiple sections with distinct headings (at least 3 sections per excerpt).\n"
            "- Never stop after one paragraph — continue until every topic is explained.\n"
            "- Each PARAGRAPH must be a detailed, multi-sentence explanation (at least 5 sentences).\n"
            "- Each BULLETS block must have 3 to 5 bullet points.\n"
            "- Do NOT summarize the excerpt in a single section.\n"
            "- Do NOT add any text before the first SECTION: line.\n\n"
            f"Excerpt:\n{chunk}\n\nExplanation:"
        )
        params = TextGenParameters(max_new_tokens=1500)
        raw_response = model.generate_text(prompt=prompt, params=params)
        raw = _extract_model_text(raw_response)
        sections = _parse_explain_sections(raw)
        all_sections.extend(sections)

    print(f"[EXPLAIN] collected {len(all_sections)} sections from "
          f"{len(chunks)} chunks", flush=True)

    return all_sections if all_sections else [
        {"heading": "Explanation", "paragraph": "No content could be extracted.", "bullets": []}
    ]


def _demo_explain(text: str) -> list[dict]:
    """
    Demo mode explanation — deterministic, no AI call.
    Splits the document into logical chunks and labels them as sections.
    """
    sentences = [s.strip() for s in text.replace("\n", " ").split(".") if len(s.strip()) > 15]
    # Group sentences into chunks of ~3, treating each group as a section
    chunk_size = 3
    chunks = [sentences[i:i + chunk_size] for i in range(0, min(len(sentences), 18), chunk_size)]
    sections = []
    for idx, chunk in enumerate(chunks, 1):
        heading = f"Section {idx}"
        para    = ". ".join(s for s in chunk if s) + ("." if chunk else "")
        bullets = [s[:100] for s in chunk]
        sections.append({"heading": heading, "paragraph": para, "bullets": bullets})
    if not sections:
        sections = [{
            "heading":   "Document Overview",
            "paragraph": "[DEMO MODE — No extractable text was found in the uploaded document.]",
            "bullets":   ["Upload a document with readable text to generate a real explanation."],
        }]
    # Prepend demo banner to the first section paragraph
    banner = (
        "[DEMO MODE — IBM watsonx.ai credentials are not configured. "
        "Set WATSONX_API_KEY, WATSONX_PROJECT_ID and WATSONX_URL for real AI explanations.]\n\n"
    )
    sections[0]["paragraph"] = banner + sections[0]["paragraph"]
    return sections


def _parse_explain_sections(raw: str) -> list[dict]:
    """
    Parse the structured output from the watsonx.ai explain prompt into a list of
    section dicts: [{ heading, paragraph, bullets }, …]
    Falls back to a single section containing the full raw text if parsing fails.
    """
    sections   = []
    current    = None

    for line in raw.splitlines():
        stripped = line.strip()
        if stripped.upper().startswith("SECTION:"):
            if current:
                sections.append(current)
            current = {"heading": stripped[8:].strip(), "paragraph": "", "bullets": []}
        elif stripped.upper().startswith("PARAGRAPH:") and current:
            current["paragraph"] = stripped[10:].strip()
        elif stripped.startswith("- ") and current:
            current["bullets"].append(stripped[2:].strip())
        elif stripped and current and not current["paragraph"]:
            current["paragraph"] = stripped

    if current:
        sections.append(current)

    if not sections:
        sections = [{"heading": "Explanation", "paragraph": raw.strip(), "bullets": []}]

    return sections


def _save_explain_file(sections: list[dict],
                       original_filename: str, session_id: str) -> str:
    """Write the explanation to translated/<session_id>_explanation.txt."""
    out_name = f"{session_id}_explanation.txt"
    out_path = os.path.join(app.config["TRANSLATED_FOLDER"], out_name)
    with open(out_path, "w", encoding="utf-8") as fh:
        fh.write(f"Explanation of: {original_filename}\n")
        fh.write("=" * 60 + "\n\n")
        for sec in sections:
            fh.write(f"{sec['heading'].upper()}\n")
            fh.write("-" * len(sec["heading"]) + "\n")
            fh.write(sec["paragraph"] + "\n")
            if sec.get("bullets"):
                fh.write("\n")
                for b in sec["bullets"]:
                    fh.write(f"  • {b}\n")
            fh.write("\n")
    return out_name


def _save_summary_file(summary: str, key_points: list[str],
                       original_filename: str, session_id: str) -> str:
    """Write summary + key points to translated/<session_id>_summary.txt."""
    out_name = f"{session_id}_summary.txt"
    out_path = os.path.join(app.config["TRANSLATED_FOLDER"], out_name)
    with open(out_path, "w", encoding="utf-8") as fh:
        fh.write(f"Summary of: {original_filename}\n")
        fh.write("-" * 60 + "\n\n")
        fh.write("SUMMARY\n\n")
        fh.write(summary + "\n\n")
        fh.write("KEY POINTS\n\n")
        for kp in key_points:
            fh.write(f"- {kp}\n")
    return out_name


def _strip_code_fences(raw: str) -> str:
    raw = raw.strip()
    if raw.startswith("```"):
        raw = raw.split("\n", 1)[1] if "\n" in raw else ""
        if raw.endswith("```"):
            raw = raw.rsplit("```", 1)[0]
    return raw.strip()


def _quiz_sentence_list(text: str) -> list[str]:
    sentences = [
        part.strip()
        for part in re.split(r"(?<=[.!?])\s+|\n+", text)
        if len(part.strip()) > 20
    ]
    if not sentences and text.strip():
        sentences = [text.strip()]
    return sentences or ["The document contains important information."]


def _quiz_keyword_pool(text: str) -> list[str]:
    pool = []
    seen = set()
    for raw in re.findall(r"[A-Za-z][A-Za-z'-]+", text):
        word = raw.strip()
        normalized = word.lower()
        if len(normalized) < 4 or normalized in _QUIZ_STOPWORDS:
            continue
        if normalized in seen:
            continue
        seen.add(normalized)
        pool.append(word)
    if not pool:
        pool = ["concept", "example", "detail", "method", "result", "process", "idea", "system"]
    return pool


def _quiz_match_case(sentence: str, word: str) -> str:
    match = re.search(re.escape(word), sentence, re.IGNORECASE)
    return match.group(0) if match else word


def _quiz_choose_target(sentence: str, keyword_pool: list[str], fallback_word: str) -> str:
    candidates = []
    for raw in re.findall(r"[A-Za-z][A-Za-z'-]+", sentence):
        word = raw.strip()
        normalized = word.lower()
        if len(normalized) < 4 or normalized in _QUIZ_STOPWORDS:
            continue
        candidates.append(word)
    if candidates:
        candidates.sort(key=lambda value: (-len(value), value.lower()))
        return candidates[0]
    if keyword_pool:
        return keyword_pool[0]
    return fallback_word


def _quiz_make_options(answer: str, keyword_pool: list[str], rng: random.Random) -> tuple[dict[str, str], str]:
    options: list[str] = []
    seen = set()

    def add_option(value: str) -> None:
        cleaned = value.strip()
        normalized = cleaned.lower()
        if cleaned and normalized not in seen:
            seen.add(normalized)
            options.append(cleaned)

    add_option(answer)
    for candidate in keyword_pool:
        if len(options) >= 4:
            break
        if candidate.lower() != answer.lower():
            add_option(candidate)

    fallback_options = [
        "concept", "detail", "example", "method", "result", "process",
        "context", "analysis", "feature", "idea", "system", "learning",
        "information", "development", "structure",
    ]
    for candidate in fallback_options:
        if len(options) >= 4:
            break
        if candidate.lower() != answer.lower():
            add_option(candidate)

    while len(options) < 4:
        add_option(f"Choice {len(options) + 1}")

    options = options[:4]
    rng.shuffle(options)

    letters = ["A", "B", "C", "D"]
    correct_letter = letters[next(index for index, value in enumerate(options) if value.lower() == answer.lower())]
    return dict(zip(letters, options)), correct_letter


def _demo_quiz(text: str, session_id: str) -> list[dict]:
    sentences = _quiz_sentence_list(text)
    keyword_pool = _quiz_keyword_pool(text)
    rng = random.Random(session_id)
    questions = []

    for index in range(5):
        sentence = sentences[index % len(sentences)]
        fallback_word = keyword_pool[index % len(keyword_pool)] if keyword_pool else "concept"
        target = _quiz_choose_target(sentence, keyword_pool, fallback_word)
        answer = _quiz_match_case(sentence, target)

        if re.search(re.escape(target), sentence, re.IGNORECASE):
            stem = re.sub(re.escape(answer), "____", sentence, count=1, flags=re.IGNORECASE)
        else:
            stem = f"{sentence.rstrip('.!?')} ____"

        options, correct_letter = _quiz_make_options(answer, keyword_pool, rng)
        questions.append({
            "id": f"q{index + 1}",
            "question": f"Which word best completes this sentence from the document?\n\n{stem}",
            "options": options,
            "answer": correct_letter,
            "explanation": f"The correct answer is {answer}.",
        })

    return questions


def _watsonx_quiz(text: str, session_id: str) -> list[dict]:
    if not WATSONX_AVAILABLE:
        raise RuntimeError("ibm-watsonx-ai SDK not installed.")
    if not WATSONX_API_KEY or not WATSONX_PROJECT_ID:
        raise RuntimeError("WATSONX_API_KEY or WATSONX_PROJECT_ID not set.")

    credentials = Credentials(url=WATSONX_URL, api_key=WATSONX_API_KEY)
    model = ModelInference(
        model_id=WATSONX_MODEL_ID,
        credentials=credentials,
        project_id=WATSONX_PROJECT_ID,
    )

    excerpt = text[:5000]
    prompt = (
        "Create exactly 5 multiple-choice questions from the document below. "
        "Return valid JSON only as an array of 5 objects. Each object must contain "
        "question, options, answer, and explanation. options must be an object with "
        "exactly the keys A, B, C, and D. answer must be one of A, B, C, or D. "
        "Do not include markdown, code fences, or extra commentary.\n\n"
        f"{excerpt}"
    )
    params = TextGenParameters(max_new_tokens=1000)
    raw = model.generate_text(prompt=prompt, params=params).strip()
    
    print(f"[DEBUG _watsonx_quiz] Raw Watsonx response:\n{raw}", flush=True)
    
    # Extract JSON array between first [ and matching ]
    start_idx = raw.find("[")
    if start_idx == -1:
        raise ValueError(f"No JSON array found in response. Raw response:\n{raw}")
    
    # Find matching closing bracket
    bracket_count = 0
    end_idx = -1
    for i in range(start_idx, len(raw)):
        if raw[i] == "[":
            bracket_count += 1
        elif raw[i] == "]":
            bracket_count -= 1
            if bracket_count == 0:
                end_idx = i + 1
                break
    
    if end_idx == -1:
        raise ValueError(f"No matching closing bracket found. Raw response:\n{raw}")
    
    extracted = raw[start_idx:end_idx].strip()
    
    # Remove any json/code fences from extracted portion
    extracted = _strip_code_fences(extracted)
    
    print(f"[DEBUG _watsonx_quiz] Extracted JSON array:\n{extracted}", flush=True)
    
    if not extracted:
        raise ValueError(f"Empty extracted JSON. Raw response:\n{raw}")

    import json as _json
    try:
        parsed = _json.loads(extracted)
    except _json.JSONDecodeError as e:
        raise ValueError(f"Invalid JSON from Watsonx: {e}. Extracted:\n{extracted}\nRaw response:\n{raw}")
    items = parsed.get("questions") if isinstance(parsed, dict) else parsed
    if not isinstance(items, list) or len(items) != 5:
        raise ValueError("Quiz payload must contain exactly 5 questions.")

    normalized = []
    for index, item in enumerate(items, 1):
        if not isinstance(item, dict):
            raise ValueError("Quiz item must be an object.")
        question = str(item.get("question", "")).strip()
        options_raw = item.get("options", {})
        if isinstance(options_raw, dict):
            options = {letter: str(options_raw.get(letter, "")).strip() for letter in ["A", "B", "C", "D"]}
        elif isinstance(options_raw, list) and len(options_raw) == 4:
            options = {letter: str(value).strip() for letter, value in zip(["A", "B", "C", "D"], options_raw)}
        else:
            raise ValueError("Quiz options must contain four choices.")

        answer = str(item.get("answer", "")).strip().upper()
        explanation = str(item.get("explanation", "")).strip()
        if not question or answer not in options or any(not value for value in options.values()):
            raise ValueError("Quiz item is missing required data.")

        normalized.append({
            "id": f"q{index}",
            "question": question,
            "options": options,
            "answer": answer,
            "explanation": explanation,
        })

    return normalized


def _build_quiz(text: str, session_id: str) -> tuple[list[dict], str]:
    mode = "watsonx"
    questions = _watsonx_quiz(text, session_id)
    return questions, mode


def _quiz_performance_message(score: int, total: int) -> str:
    percentage = round((score / total) * 100) if total else 0
    if percentage == 100:
        return "Excellent work. You answered every question correctly."
    if percentage >= 80:
        return "Strong result. You have a solid grasp of the document."
    if percentage >= 60:
        return "Good effort. Review the missed questions once more."
    return "Keep studying the document and try again."


def _quiz_history_preview(score: int, total: int, percentage: int) -> str:
    return f"Quiz completed: {score}/{total} correct ({percentage}%)."


def _record_quiz_history(session_id: str, doc: dict, mode: str, quiz_id: str,
                         score: int, total: int, percentage: int) -> None:
    history_store.append({
        "id": quiz_id,
        "type": "quiz",
        "session_id": session_id,
        "filename": doc["filename"],
        "target_language": "",
        "word_count": doc.get("word_count", 0),
        "mode": mode,
        "score": score,
        "total_questions": total,
        "percentage": percentage,
        "timestamp": datetime.datetime.utcnow().strftime("%Y-%m-%d %H:%M UTC"),
        "preview": _quiz_history_preview(score, total, percentage),
    })


def _ask_stopwords() -> set[str]:
    return {
        "about", "after", "again", "also", "because", "been", "before", "being",
        "between", "both", "could", "during", "each", "from", "have", "having",
        "here", "into", "many", "more", "most", "only", "other", "some",
        "such", "than", "that", "them", "then", "there", "these", "they",
        "this", "those", "through", "under", "using", "very", "when", "where",
        "which", "while", "with", "without", "would", "your", "their", "will",
        "shall", "onto", "over", "under", "what", "who", "whom", "why", "how",
        "the", "and", "for", "are", "was", "were", "has", "had", "have",
        "not", "but", "can", "may", "our", "out", "use", "used", "useful",
        "document", "question", "answer", "please", "tell", "show", "give",
    }


def _ask_normalize_query(text: str) -> list[str]:
    words = []
    stopwords = _ask_stopwords()
    for raw in re.findall(r"[A-Za-z][A-Za-z'-]+", text):
        word = raw.lower().strip()
        if len(word) < 3 or word in stopwords:
            continue
        if word not in words:
            words.append(word)
    return words


def _ask_sentence_chunks(text: str) -> list[str]:
    chunks = []
    for part in re.split(r"\n+|(?<=[.!?])\s+", text):
        cleaned = part.strip()
        if len(cleaned) > 25:
            chunks.append(cleaned)
    return chunks


def _ask_retrieve_context(text: str, question: str, conversation: list[dict] | None = None) -> list[str]:
    query_terms = _ask_normalize_query(question)
    if conversation:
        for item in conversation[-6:]:
            query_terms.extend(_ask_normalize_query(str(item.get("content", ""))))

    query_terms = list(dict.fromkeys(query_terms))
    chunks = _ask_sentence_chunks(text)
    if not chunks:
        return [text[:1000].strip()] if text.strip() else []

    scored: list[tuple[int, str]] = []
    for chunk in chunks:
        lower_chunk = chunk.lower()
        score = 0
        for term in query_terms:
            if term in lower_chunk:
                score += 2 if len(term) > 4 else 1
        if score:
            scored.append((score, chunk))

    if not scored:
        return chunks[:4]

    scored.sort(key=lambda item: (-item[0], len(item[1])))
    selected = []
    for _, chunk in scored:
        if chunk not in selected:
            selected.append(chunk)
        if len(selected) >= 6:
            break
    return selected


def _ask_build_context_block(snippets: list[str]) -> str:
    if not snippets:
        return ""
    return "\n\n".join(f"- {snippet}" for snippet in snippets)


def _ask_compose_prompt(question: str, document_text: str, conversation: list[dict] | None = None) -> str:
    snippets = _ask_retrieve_context(document_text, question, conversation)
    context_block = _ask_build_context_block(snippets)

    recent_turns = []
    if conversation:
        for item in conversation[-6:]:
            role = str(item.get("role", "user")).lower()
            content = str(item.get("content", "")).strip()
            if not content:
                continue
            label = "User" if role == "user" else "Assistant"
            recent_turns.append(f"{label}: {content}")

    conversation_block = "\n".join(recent_turns).strip()
    prompt = (
        "You are a document question-answering assistant. Answer using only the provided document context. "
        "If the answer is not present in the document, say that the document does not contain enough information. "
        "Keep the response concise, clear, and helpful.\n\n"
        f"Document context:\n{context_block if context_block else document_text[:4000]}\n\n"
    )
    if conversation_block:
        prompt += f"Conversation so far:\n{conversation_block}\n\n"
    prompt += f"Question: {question}\nAnswer:"
    return prompt


def _watsonx_ask_answer(question: str, document_text: str, conversation: list[dict] | None = None) -> tuple[str, list[str]]:
    if not WATSONX_AVAILABLE:
        raise RuntimeError("ibm-watsonx-ai SDK not installed.")
    if not WATSONX_API_KEY or not WATSONX_PROJECT_ID:
        raise RuntimeError("WATSONX_API_KEY or WATSONX_PROJECT_ID not set.")

    credentials = Credentials(url=WATSONX_URL, api_key=WATSONX_API_KEY)
    model = ModelInference(
        model_id=WATSONX_MODEL_ID,
        credentials=credentials,
        project_id=WATSONX_PROJECT_ID,
    )

    snippets = _ask_retrieve_context(document_text, question, conversation)
    prompt = _ask_compose_prompt(question, document_text, conversation)
    params = TextGenParameters(max_new_tokens=500)
    answer = str(model.generate_text(prompt=prompt, params=params)).strip()
    if not answer:
        raise ValueError("Empty answer returned from model.")
    return answer, snippets


def _demo_ask_answer(question: str, document_text: str, conversation: list[dict] | None = None) -> tuple[str, list[str]]:
    snippets = _ask_retrieve_context(document_text, question, conversation)
    if not snippets:
        return (
            "The uploaded document does not contain enough information to answer that question.",
            [],
        )

    lead = snippets[0]
    extra = snippets[1:3]
    if extra:
        snippet_text = " ".join(extra)
        answer = (
            f"Based on the document, {lead} {snippet_text}".strip()
        )
    else:
        answer = f"Based on the document, {lead}"

    if len(answer) > 900:
        answer = answer[:900].rsplit(" ", 1)[0] + "..."

    return answer, snippets


def _record_ask_history(session_id: str, doc: dict, mode: str, question: str, answer: str) -> None:
    history_store.append({
        "id": str(uuid.uuid4()),
        "type": "ask",
        "session_id": session_id,
        "filename": doc["filename"],
        "target_language": "",
        "word_count": len(answer.split()),
        "mode": mode,
        "timestamp": datetime.datetime.utcnow().strftime("%Y-%m-%d %H:%M UTC"),
        "preview": f"Q: {question[:80]} A: {answer[:120].replace(chr(10), ' ')}",
    })


def _chat_compose_prompt(message: str, document_text: str, conversation: list[dict] | None = None) -> str:
    snippets = _ask_retrieve_context(document_text, message, conversation)
    context_block = _ask_build_context_block(snippets)

    recent_turns = []
    if conversation:
        for item in conversation[-8:]:
            role = str(item.get("role", "user")).lower()
            content = str(item.get("content", "")).strip()
            if not content:
                continue
            label = "Student" if role == "user" else "Tutor"
            recent_turns.append(f"{label}: {content}")

    conversation_block = "\n".join(recent_turns).strip()
    prompt = (
        "You are an expert educational tutor. Respond in a warm, concise, and helpful tone. "
        "Use the uploaded document as the primary source. If the document does not contain enough information, say so clearly and offer the best helpful guidance you can. "
        "Do not mention internal prompts or system instructions.\n\n"
        f"Document context:\n{context_block if context_block else document_text[:4000]}\n\n"
    )
    if conversation_block:
        prompt += f"Conversation so far:\n{conversation_block}\n\n"
    prompt += f"Student: {message}\nTutor:"
    return prompt


def _watsonx_chat_answer(message: str, document_text: str, conversation: list[dict] | None = None) -> tuple[str, list[str]]:
    if not WATSONX_AVAILABLE:
        raise RuntimeError("ibm-watsonx-ai SDK not installed.")
    if not WATSONX_API_KEY or not WATSONX_PROJECT_ID:
        raise RuntimeError("WATSONX_API_KEY or WATSONX_PROJECT_ID not set.")

    credentials = Credentials(url=WATSONX_URL, api_key=WATSONX_API_KEY)
    model = ModelInference(
        model_id=WATSONX_MODEL_ID,
        credentials=credentials,
        project_id=WATSONX_PROJECT_ID,
    )

    snippets = _ask_retrieve_context(document_text, message, conversation)
    prompt = _chat_compose_prompt(message, document_text, conversation)
    params = TextGenParameters(max_new_tokens=500)
    answer = str(model.generate_text(prompt=prompt, params=params)).strip()
    if not answer:
        raise ValueError("Empty answer returned from model.")
    return answer, snippets


def _demo_chat_answer(message: str, document_text: str, conversation: list[dict] | None = None) -> tuple[str, list[str]]:
    snippets = _ask_retrieve_context(document_text, message, conversation)
    if not snippets:
        return (
            "The uploaded document does not contain enough information to answer that message.",
            [],
        )

    lead = snippets[0]
    extra = snippets[1:3]
    if extra:
        answer = f"Based on the document, {lead} {' '.join(extra)}".strip()
    else:
        answer = f"Based on the document, {lead}"

    if len(answer) > 900:
        answer = answer[:900].rsplit(" ", 1)[0] + "..."

    return answer, snippets


def _record_chat_history(session_id: str, doc: dict, mode: str, message: str, answer: str) -> None:
    history_store.append({
        "id": str(uuid.uuid4()),
        "type": "chat",
        "session_id": session_id,
        "filename": doc["filename"],
        "target_language": "",
        "word_count": len(answer.split()),
        "mode": mode,
        "timestamp": datetime.datetime.utcnow().strftime("%Y-%m-%d %H:%M UTC"),
        "preview": f"Q: {message[:80]} A: {answer[:120].replace(chr(10), ' ')}",
    })


def save_translated_file(translated_text: str, original_filename: str,
                         target_language: str, session_id: str) -> str:
    """
    Write the translated text to translated/<session_id>_<lang>.txt
    and return the saved filename.
    """
    safe_lang  = target_language.replace(" ", "_")
    out_name   = f"{session_id}_{safe_lang}.txt"
    out_path   = os.path.join(app.config["TRANSLATED_FOLDER"], out_name)
    with open(out_path, "w", encoding="utf-8") as fh:
        fh.write(f"Translated to: {target_language}\n")
        fh.write(f"Source file:   {original_filename}\n")
        fh.write("-" * 60 + "\n\n")
        fh.write(translated_text)
    return out_name


# ── Page routes ────────────────────────────────────────────────────────────

@app.route("/")
def index():
    return render_template("index.html")


@app.route("/translate")
def translate():
    return render_template("index.html", active="translate")


@app.route("/summarize")
def summarize():
    return render_template("index.html", active="summarize")


@app.route("/explain")
def explain():
    return render_template("index.html", active="explain")


@app.route("/quiz")
def quiz():
    return render_template("index.html", active="quiz")


@app.route("/ask")
def ask():
    return render_template("index.html", active="ask")


@app.route("/chat")
def chat():
    return render_template("index.html", active="chat")


@app.route("/history")
def history():
    return render_template("index.html", active="history")


@app.route("/api/quiz", methods=["POST"])
def api_quiz():
    """
    Generate a quiz from a previously uploaded document or evaluate a submitted quiz.

    Generation request JSON: { session_id }
    Submission request JSON: { session_id, quiz_id, answers }
    """
    body = request.get_json(silent=True) or {}
    session_id = body.get("session_id", "").strip()
    quiz_id = body.get("quiz_id", "").strip()
    answers = body.get("answers", {})

    if not session_id:
        return jsonify({"success": False, "error": "session_id is required."}), 400

    doc = document_store.get(session_id)
    if doc is None:
        return jsonify({
            "success": False,
            "error": "Session not found. Please upload a document first.",
        }), 404

    source_text = doc.get("text", "")
    if not source_text.strip():
        return jsonify({
            "success": False,
            "error": "The uploaded document contains no extractable text.",
        }), 422

    if quiz_id:
        stored_quiz = quiz_store.get(quiz_id)
        if stored_quiz is None or stored_quiz.get("session_id") != session_id:
            return jsonify({
                "success": False,
                "error": "Quiz session not found. Please generate the quiz again.",
            }), 404

        if not isinstance(answers, dict):
            return jsonify({"success": False, "error": "answers must be an object."}), 400

        evaluations = []
        score = 0
        for question in stored_quiz["questions"]:
            selected = str(answers.get(question["id"], "")).strip().upper()
            correct = question["answer"]
            is_correct = selected == correct
            if is_correct:
                score += 1
            evaluations.append({
                "question_id": question["id"],
                "selected": selected,
                "correct_answer": correct,
                "is_correct": is_correct,
            })

        total = len(stored_quiz["questions"])
        percentage = round((score / total) * 100) if total else 0
        performance_message = _quiz_performance_message(score, total)

        _record_quiz_history(session_id, doc, stored_quiz["mode"], quiz_id, score, total, percentage)

        return jsonify({
            "success": True,
            "quiz_id": quiz_id,
            "score": score,
            "total_questions": total,
            "percentage": percentage,
            "performance_message": performance_message,
            "results": evaluations,
            "mode": stored_quiz["mode"],
            "filename": doc["filename"],
        })

    questions, mode = _build_quiz(source_text, session_id)
    generated_quiz_id = str(uuid.uuid4())
    quiz_store[generated_quiz_id] = {
        "session_id": session_id,
        "filename": doc["filename"],
        "questions": questions,
        "mode": mode,
        "created_at": datetime.datetime.utcnow().strftime("%Y-%m-%d %H:%M UTC"),
    }

    safe_questions = [
        {
            "id": question["id"],
            "question": question["question"],
            "options": question["options"],
        }
        for question in questions
    ]

    return jsonify({
        "success": True,
        "quiz_id": generated_quiz_id,
        "questions": safe_questions,
        "question_count": len(safe_questions),
        "mode": mode,
        "filename": doc["filename"],
        "word_count": doc.get("word_count", 0),
    })


@app.route("/api/ask", methods=["POST"])
def api_ask():
    """
    Answer a question about the currently uploaded document.

    Request JSON: { session_id, question, conversation? }
    Response JSON: { success, answer, mode, filename, question, sources }
    """
    body = request.get_json(silent=True) or {}
    session_id = body.get("session_id", "").strip()
    question = body.get("question", "").strip()
    conversation = body.get("conversation", [])

    if not session_id:
        return jsonify({"success": False, "error": "session_id is required."}), 400

    doc = document_store.get(session_id)
    if doc is None:
        return jsonify({
            "success": False,
            "error": "Session not found. Please upload a document first.",
        }), 404

    if not question:
        return jsonify({"success": False, "error": "question is required."}), 400

    if not isinstance(conversation, list):
        conversation = []

    source_text = doc.get("text", "")
    if not source_text.strip():
        return jsonify({
            "success": False,
            "error": "The uploaded document contains no extractable text.",
        }), 422

    mode = "watsonx"
    answer, sources = _watsonx_ask_answer(question, source_text, conversation)

    _record_ask_history(session_id, doc, mode, question, answer)

    return jsonify({
        "success": True,
        "answer": answer,
        "mode": mode,
        "filename": doc["filename"],
        "question": question,
        "sources": sources,
    })


@app.route("/api/chat", methods=["POST"])
def api_chat():
    """
    Continue a document-grounded chat conversation using the uploaded file associated with the current session.

    Request JSON: { session_id, message, conversation? }
    Response JSON: { success, answer, mode, filename, message, sources }
    """
    body = request.get_json(silent=True) or {}
    session_id = body.get("session_id", "").strip()
    message = body.get("message", "").strip()
    conversation = body.get("conversation", [])

    if not session_id:
        return jsonify({"success": False, "error": "session_id is required."}), 400

    doc = document_store.get(session_id)
    if doc is None:
        return jsonify({
            "success": False,
            "error": "Session not found. Please upload a document first.",
        }), 404

    if not message:
        return jsonify({"success": False, "error": "message is required."}), 400

    if not isinstance(conversation, list):
        conversation = []

    source_text = doc.get("text", "")
    if not source_text.strip():
        return jsonify({
            "success": False,
            "error": "The uploaded document contains no extractable text.",
        }), 422

    mode = "watsonx"
    answer, sources = _watsonx_chat_answer(message, source_text, conversation)

    _record_chat_history(session_id, doc, mode, message, answer)

    return jsonify({
        "success": True,
        "answer": answer,
        "mode": mode,
        "filename": doc["filename"],
        "message": message,
        "sources": sources,
    })


# ── API routes ─────────────────────────────────────────────────────────────

@app.route("/api/upload", methods=["POST"])
def api_upload():
    """
    Accept a multipart file upload, validate, parse and store it.
    Returns JSON: { success, filename, session_id, file_type, pages, word_count }
    """
    if "file" not in request.files:
        return jsonify({"success": False, "error": "No file field in request."}), 400

    file = request.files["file"]

    if file.filename == "":
        return jsonify({"success": False, "error": "No file selected."}), 400

    if not allowed_extension(file.filename):
        return jsonify({
            "success": False,
            "error": "Unsupported file type. Allowed: PDF, DOCX, PPTX, TXT.",
        }), 415

    ext        = file.filename.rsplit(".", 1)[1].lower()
    session_id = str(uuid.uuid4())
    safe_name  = f"{session_id}.{ext}"
    save_path  = os.path.join(app.config["UPLOAD_FOLDER"], safe_name)

    file.save(save_path)

    try:
        text, pages = extract_text(save_path, ext)
    except Exception as exc:
        os.remove(save_path)
        return jsonify({"success": False, "error": f"Could not parse file: {exc}"}), 422

    word_count = len(text.split())

    document_store[session_id] = {
        "text":       text,
        "filename":   file.filename,
        "file_type":  ext.upper(),
        "pages":      pages,
        "word_count": word_count,
    }

    # Persist immediately so a Flask reload doesn't lose this session.
    _flush_document_store()

    return jsonify({
        "success":    True,
        "filename":   file.filename,
        "session_id": session_id,
        "file_type":  ext.upper(),
        "pages":      pages,
        "word_count": word_count,
    })


@app.route("/api/translate", methods=["POST"])
def api_translate():
    """
    Translate a previously uploaded document.

    Request JSON: { session_id, target_language }
    Response JSON:
        { success, translated_text, word_count, target_language,
          filename, mode, saved_file }
    """
    body = request.get_json(silent=True) or {}

    session_id      = body.get("session_id", "").strip()
    target_language = body.get("target_language", "").strip()

    if not session_id:
        return jsonify({"success": False, "error": "session_id is required."}), 400

    if target_language not in SUPPORTED_LANGUAGES:
        return jsonify({
            "success": False,
            "error": f"Unsupported language. Choose from: {', '.join(SUPPORTED_LANGUAGES)}.",
        }), 400

    doc = document_store.get(session_id)
    if doc is None:
        return jsonify({
            "success": False,
            "error": "Session not found. Please upload a document first.",
        }), 404

    source_text = doc["text"]
    if not source_text.strip():
        return jsonify({
            "success": False,
            "error": "The uploaded document contains no extractable text.",
        }), 422

    # ── Use watsonx.ai ─────────────────────────────────────────────────────────
    mode = "watsonx"
    translated_text = _watsonx_translate(source_text, target_language)

    word_count = len(translated_text.split())

    # ── Persist to translated/ folder ─────────────────────────────────────
    saved_file = save_translated_file(
        translated_text, doc["filename"], target_language, session_id
    )

    # ── Record in history ──────────────────────────────────────────────────
    history_store.append({
        "id":              str(uuid.uuid4()),
        "type":            "translation",
        "session_id":      session_id,
        "filename":        doc["filename"],
        "target_language": target_language,
        "word_count":      word_count,
        "mode":            mode,
        "saved_file":      saved_file,
        "timestamp":       datetime.datetime.utcnow().strftime("%Y-%m-%d %H:%M UTC"),
        "preview":         translated_text[:120].replace("\n", " "),
    })

    return jsonify({
        "success":         True,
        "translated_text": translated_text,
        "word_count":      word_count,
        "target_language": target_language,
        "filename":        doc["filename"],
        "mode":            mode,
        "saved_file":      saved_file,
    })


@app.route("/api/summarize", methods=["POST"])
def api_summarize():
    """
    Summarise a previously uploaded document.

    Request JSON : { session_id }
    Response JSON: { success, summary, key_points, word_count,
                     filename, mode, saved_file }
    """
    body       = request.get_json(silent=True) or {}
    session_id = body.get("session_id", "").strip()

    if not session_id:
        return jsonify({"success": False, "error": "session_id is required."}), 400

    doc = document_store.get(session_id)
    if doc is None:
        return jsonify({
            "success": False,
            "error": "Session not found. Please upload a document first.",
        }), 404

    source_text = doc["text"]
    if not source_text.strip():
        return jsonify({
            "success": False,
            "error": "The uploaded document contains no extractable text.",
        }), 422

    # ── Use watsonx.ai ─────────────────────────────────────────────────────────
    mode = "watsonx"
    summary, key_points = _watsonx_summarize(source_text)

    word_count = len(summary.split())

    # ── Persist output ─────────────────────────────────────────────────────
    saved_file = _save_summary_file(summary, key_points, doc["filename"], session_id)

    # ── Record in history ──────────────────────────────────────────────────
    history_store.append({
        "id":              str(uuid.uuid4()),
        "type":            "summary",
        "session_id":      session_id,
        "filename":        doc["filename"],
        "target_language": "",
        "word_count":      word_count,
        "mode":            mode,
        "saved_file":      saved_file,
        "timestamp":       datetime.datetime.utcnow().strftime("%Y-%m-%d %H:%M UTC"),
        "preview":         summary[:120].replace("\n", " "),
    })

    return jsonify({
        "success":    True,
        "summary":    summary,
        "key_points": key_points,
        "word_count": word_count,
        "filename":   doc["filename"],
        "mode":       mode,
        "saved_file": saved_file,
    })


@app.route("/api/explain", methods=["POST"])
def api_explain():
    """
    Generate a chapter-wise explanation for a previously uploaded document.

    Request JSON : { session_id }
    Response JSON: { success, sections, word_count, filename, mode, saved_file }
      sections = [{ heading, paragraph, bullets }, …]
    """
    body       = request.get_json(silent=True) or {}
    session_id = body.get("session_id", "").strip()

    if not session_id:
        return jsonify({"success": False, "error": "session_id is required."}), 400

    doc = document_store.get(session_id)
    if doc is None:
        return jsonify({
            "success": False,
            "error": "Session not found. Please upload a document first.",
        }), 404

    source_text = doc["text"]
    if not source_text.strip():
        return jsonify({
            "success": False,
            "error": "The uploaded document contains no extractable text.",
        }), 422

    # ── Use watsonx.ai ─────────────────────────────────────────────────────────
    mode = "watsonx"
    sections = _watsonx_explain(source_text)

    # Flat text for word-count and preview
    flat_text  = " ".join(
        s["paragraph"] + " " + " ".join(s.get("bullets", []))
        for s in sections
    )
    word_count = len(flat_text.split())
    preview    = sections[0]["paragraph"][:120].replace("\n", " ") if sections else ""

    saved_file = _save_explain_file(sections, doc["filename"], session_id)

    history_store.append({
        "id":              str(uuid.uuid4()),
        "type":            "explanation",
        "session_id":      session_id,
        "filename":        doc["filename"],
        "target_language": "",
        "word_count":      word_count,
        "mode":            mode,
        "saved_file":      saved_file,
        "timestamp":       datetime.datetime.utcnow().strftime("%Y-%m-%d %H:%M UTC"),
        "preview":         preview,
    })

    return jsonify({
        "success":    True,
        "sections":   sections,
        "word_count": word_count,
        "filename":   doc["filename"],
        "mode":       mode,
        "saved_file": saved_file,
    })


@app.route("/api/history", methods=["GET"])
def api_history():
    """Return the full history list, newest first."""
    return jsonify({
        "success": True,
        "history": list(reversed(history_store)),
    })


# ── Entry point ────────────────────────────────────────────────────────────

if __name__ == "__main__":
    app.run(debug=True)
